import os
import io
import sys
import csv
import glob
import subprocess
from tqdm import tqdm
import numpy as np
import nibabel as nib
from PIL import Image
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# alpha = 0.3
# alpha = 0.5
alpha = 0.75
cmap_name = 'tab20b'
max_outdim = 500
# max_outdim = 256

lower_thresh = 0.0
upper_thresh = 0.4
def normalize(img):
    _min = img.min()
    _max = img.max()
    _range = _max - _min
    img[img < _range * lower_thresh] = _range * lower_thresh
    img[img > _range * upper_thresh] = _range * upper_thresh
    _min = img.min()
    _max = img.max()
    img = (img - _min) / (_max - _min)
    return (255 * img).astype(np.uint8)

def get_slice(img, center, slice_loc):
    if slice_loc == 0:
        return img[center[0], :, :]
    elif slice_loc == 1:
        return img[:, center[1], :]
    else:
        return img[:, :, center[2]]

extracranial_values = [165, 258, 259]
# seg_vals = [0, 2, 3, 4, 5, 7, 8, 10, 11, 12, 13, 14, 15, 16, 17, 18, 24, 26, 28, 30, 31, 41, 42, 43, 44, 46, 47, 49, 50, 51, 52, 53, 54, 58, 60, 62, 63, 72, 77, 80, 85, 165, 258, 259] - extracranial_values
seg_vals = [0, 2, 3, 4, 5, 7, 8, 10, 11, 12, 13, 14, 15, 16, 17, 18, 24, 26, 28, 30, 31, 41, 42, 43, 44, 46, 47, 49, 50, 51, 52, 53, 54, 58, 60, 62, 63, 72, 77, 80, 85]

def apply_colormap(img):
    # seg_values = np.unique(img)
    seg_mapping = {}
    # for idx, val in enumerate(seg_values):
    for idx, val in enumerate(seg_vals):
        img[img == val] = idx
    # cm = plt.get_cmap(cmap_name, lut=len(seg_values))
    cm = plt.get_cmap(cmap_name, lut=len(seg_vals))
    mapped = cm(img.astype(np.uint8))
    return (255 * mapped).astype(np.uint8)

slide_height = 5.625
slide_width = 10

img_dim = 2.5

def clean_mask(mask):
    for v in extracranial_values:
        mask[mask == v] = 0
    return mask

def pad_img(img):
    old_size = img.size
    new_size = (max_outdim, max_outdim)
    new_img = Image.new("RGBA", new_size)
    new_img.paste(img, (int((new_size[0]-old_size[0])/2),
                          int((new_size[1]-old_size[1])/2)))
    return new_img

rot_nums = {
    # 1: 90,
    # 2: 90,
}

def resize_and_rotate(img, slice_id):
    scale_factor = max_outdim / max(img.width, img.height)
    outsize = (int(img.width * scale_factor), int(img.height * scale_factor))
    # img = img.resize(outsize, Image.NEAREST)
    img = img.resize((max_outdim, max_outdim), Image.NEAREST)
    img = pad_img(img)
    return img.rotate(rot_nums.get(slice_id, 0))

bb_pad = 0
def get_bounding_box(mask):
    r = np.any(mask, axis=(1, 2))
    c = np.any(mask, axis=(0, 2))
    z = np.any(mask, axis=(0, 1))
    rmin, rmax = np.where(r)[0][[0, -1]]
    cmin, cmax = np.where(c)[0][[0, -1]]
    zmin, zmax = np.where(z)[0][[0, -1]]
    return (
        rmin - bb_pad,
        rmax + bb_pad,
        cmin - bb_pad,
        cmax + bb_pad,
        zmin - bb_pad,
        zmax + bb_pad
    )

def crop(mask, img):
    xmin, xmax, ymin, ymax, zmin, zmax = get_bounding_box(mask)
    return img[xmin:xmax, ymin:ymax, zmin:zmax]

locations = {
    0: {
        'left': Inches(2.62),
        'width': Inches(img_dim),
        'height': Inches(img_dim),
    },
    1: {
        'left': Inches(5.22),
        'width': Inches(img_dim),
        'height': Inches(img_dim),
    },
    2: {
        'left': Inches(7.41),
        'width': Inches(img_dim),
        'height': Inches(img_dim),
    }
}

def create_slides(base_path, original_volume, compare_volume, original_samseg, compare_samseg, subjs, pptx_name, normalize_original=True, normalize_compare=True):
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    blank_slide_layout = prs.slide_layouts[6]

    for idx, subj in tqdm(enumerate(subjs[:5]), total=len(subjs)):

        original_nii_path = os.path.join(base_path, subj, original_volume)
        compare_nii_path = os.path.join(base_path, subj, compare_volume)

        original_mask_path = os.path.join(base_path, subj, 'samseg', original_samseg, 'seg.mgz')
        compare_mask_path = os.path.join(base_path, subj, 'samseg', compare_samseg, 'seg.mgz')

        exists = True
        for p in [original_nii_path, compare_nii_path, compare_mask_path]:
            if not os.path.exists(p):
                print('[{}] {} does not exist'.format(subj, p))
                exists = False
        if not exists:
            continue

        slide = prs.slides.add_slide(blank_slide_layout)
        txBox = slide.shapes.add_textbox(Inches(3.7), Inches(5.21), Inches(2.6), Inches(0.4))
        txBox.text_frame.text = subj
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        original = nib.load(original_nii_path)
        compare = nib.load(compare_nii_path)
        original_mask = nib.load(original_mask_path)
        compare_mask = nib.load(compare_mask_path)

        original_data = original.get_fdata()
        if len(original_data.shape) == 4 and original_data.shape[3] == 1:
            original_data = original_data[:, :, :, 0]
        compare_data = compare.get_fdata()
        if len(compare_data.shape) == 4 and compare_data.shape[3] == 1:
            compare_data = compare_data[:, :, :, 0]

        original_mask_data = original_mask.get_fdata()
        compare_mask_data = compare_mask.get_fdata()

        original_data = crop(original_mask_data, original_data)
        compare_data = crop(original_mask_data, compare_data)

        compare_mask_data = crop(original_mask_data, compare_mask_data)
        original_mask_data = crop(original_mask_data, original_mask_data)


        original_center = [int(v / 2) for v in original_data.shape]
        compare_center = [int(v / 2) for v in compare_data.shape]
        
        print(original_center, compare_center)

        # for slice_id, top in [(2, 0), (1, 2.5)]:
        for slice_id, top in [(2, 0), (0, 2.5)]:
            original_slice = get_slice(original_data, original_center, slice_id)
            if normalize_original:
                original_slice = normalize(original_slice)
            compare_slice = get_slice(compare_data, compare_center, slice_id)
            if normalize_compare:
                compare_slice = normalize(compare_slice)


            original_mask_slice = apply_colormap(clean_mask(get_slice(original_mask_data, original_center, slice_id)))
            compare_mask_slice = apply_colormap(clean_mask(get_slice(compare_mask_data, compare_center, slice_id)))

            original_img = Image.fromarray(original_slice, 'L').convert('RGBA')
            compare_img = Image.fromarray(compare_slice, 'L').convert('RGBA')

            original_mask_img = Image.fromarray(original_mask_slice, 'RGBA')
            original_mask_img.putalpha(int(256 * alpha))

            compare_mask_img = Image.fromarray(compare_mask_slice, 'RGBA')
            compare_mask_img.putalpha(int(256 * alpha))

            original_overlaid_img = resize_and_rotate(Image.alpha_composite(original_img, original_mask_img), slice_id)
            compare_overlaid_img = resize_and_rotate(Image.alpha_composite(compare_img, compare_mask_img), slice_id)
            
            original_img = resize_and_rotate(original_img, slice_id)
            compare_img = resize_and_rotate(compare_img, slice_id)

            # original_overlaid_img = original_overlaid_img.transpose(Image.FLIP_LEFT_RIGHT)
            # original_img = original_img.transpose(Image.FLIP_LEFT_RIGHT)

            placements = [
                (original_img, 0),
                (original_overlaid_img, 2.5),
                (compare_overlaid_img, 5),
                (compare_img, 7.5)
            ]

            for img, left in placements:
                with io.BytesIO() as img_bytes:
                    img.save(img_bytes, format='png')
                    slide.shapes.add_picture(
                        img_bytes,
                        Inches(left),
                        Inches(top),
                        width=Inches(img_dim),
                        height=Inches(img_dim))

        # txBox = slide.shapes.add_textbox(Inches(5.55), Inches(2.78), Inches(0.34), Inches(0.4))
        # p = txBox.text_frame.add_paragraph()
        # txBox.text_frame.text = 'R'
        # txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # txBox.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # txBox = slide.shapes.add_textbox(Inches(7.73), Inches(2.78), Inches(0.34), Inches(0.4))
        # p = txBox.text_frame.add_paragraph()
        # txBox.text_frame.text = 'R'
        # txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # txBox.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    prs.save('pptxs/{}.pptx'.format(pptx_name))


base_path = '/data/mradovan/7T_WMn_3T_CSFn_pairs/'
subjs = sorted([s.split('/')[-1] for s in glob.glob(os.path.join(base_path, '*'))])

original = 'CSFn'
compare = 'CSFnSB_wmn_invwarped'
normalize_compare = False
name = '{}_vs_{}'.format(original, compare)

original_volume = '{}.nii.gz'.format(original)
original_samseg = original

compare_volume = '{}.nii.gz'.format(compare)
compare_samseg = compare

create_slides(base_path, original_volume, compare_volume, original_samseg, compare_samseg, subjs, name, normalize_compare=normalize_compare)
